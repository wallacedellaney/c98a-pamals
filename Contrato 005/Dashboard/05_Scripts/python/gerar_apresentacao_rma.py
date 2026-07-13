"""
Gera a apresentação RMA (PowerPoint) de um mês de referência — teste pedido
pelo Wallace em 2026-07-11 ("Junho V2 TESTE"). Ver
00_Instrucoes/apresentacao_rma.md.

Reaproveita `04_Relatorios/RMA_referencia.pptx` (cópia da apresentação real
de junho/2026, baixada do Drive) só como BASE de layout/estilo — a intenção
(Wallace, 2026-07-11) é puxar tudo da nossa própria plataforma, com exceção
da Utilização (só existe como PDF externo, fora do nosso alcance). Descoberta
em 2026-07-12: os slides de Empréstimos/Reparáveis na referência NÃO são
imagem "oficial" nenhuma — são capturas de tela do NOSSO PRÓPRIO site que o
Wallace colou manualmente lá todo mês; por isso dá pra reconstruir com dado
vivo, no lugar de pedir pra ele printar a tela de novo.

Escopo desta versão (teste, não a versão final):
- GERADO automaticamente:
  - Slide 1 (título) e slide 2 (objetivo): mês trocado pro mês escolhido,
    resto do texto intocado.
  - Slides 3 e 10 (Tabela 1.2 — matriz aeronave x dia + MMAM/P/PMAX/IFD):
    reconstruídos como tabela nativa a partir do Cômputo Mensal, com as
    mesmas cores (verde/amarelo) e as mesmas fórmulas da planilha oficial
    (conferidas célula a célula, ver calcular_computo_mensal.py):
    MMAM = média das médias diárias; P = soma de todo 0/1 da matriz;
    PMAX = quantidade de células preenchidas; IFD = MMAM arredondado PRA
    CIMA pra faixa de 0,05 mais próxima (>=95% -> 1,00; >=90% -> 0,95;
    >=85% -> 0,90; >=80% -> 0,85; >=75% -> 0,80; <75% -> 0,75).
  - Pagamentos: tabela nativa (base_pagamentos_tratada.xlsx) + resumo do
    contrato, no lugar da imagem da referência.
  - Empréstimos: slide de estatísticas (3 gráficos, mesmos de
    emprestimos.py) + tabela com todos os itens com pedido de envio no mês
    de referência.
  - Reparáveis: slide de estatísticas (gráfico "por condição", 2 visões
    lado a lado) + 2 slides de tabela (nº total + amostra) — todas em
    aberto, e só as que estão com a empresa (VEE ONE) ou nas terceirizadas
    (mesmo filtro "Onde se encontra" que o Wallace já usa no site).
  - Atrasos: 2 slides novos — "Situação Atual" (3 métricas + as
    emergências em aberto) e "Entregas no Mês de Referência" (resumo +
    rosca No prazo/Atrasado + amostra), mesmas 2 seções de `_atrasos()`
    em fechamento_mensal.py.
  - 1 slide por aeronave que negativou (tabela da(s) emergência(s)
    responsável(is) + linha 1/0 do mês).
- REMOVIDO: os slides-imagem dos antigos dashboards de Empréstimos e
  Atrasos (capturas de tela do site) — viraram redundantes com os slides
  de estatísticas nativos.
- MANTIDO da referência, sem regenerar ainda: utilização (imagem — fonte é
  um PDF externo, fica assim de propósito), tabela de faturamento e notas
  fiscais — ver 00_Instrucoes/apresentacao_rma.md pra o que falta.
- SEM "brand" de slide — tentado (fundo escuro + âmbar, identidade do
  site) e revertido a pedido do Wallace em 2026-07-12 ("ficou péssimo").
  Os gráficos novos usam a paleta de cor de DADO do site (âmbar/ciano/
  verde/vermelho), que é diferente de "brand" (isso é só a cor normal de
  gráfico, não fundo/tema da lâmina).

2026-07-13: script recriado (o arquivo tinha sido removido da pasta a
pedido do Wallace, "vc tirou a parte do slide?" / "sim, é para ter" —
recuperado a partir do bytecode ainda em cache (__pycache__/*.pyc) +
`RMA_referencia.pptx` baixado de novo do Drive (pasta de Junho/2026).
"""

import json
from datetime import date, datetime
from pathlib import Path

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from common import DADOS_TRATADOS, RELATORIOS
from gerar_ata_reuniao import _localizar_pasta_mes, _baixar_rma_em_andamento, extrair_indicadores_rma

REFERENCIA = RELATORIOS / "RMA_referencia.pptx"

MESES_PT_MAIUSCULO = [
    "JANEIRO", "FEVEREIRO", "MARÇO", "ABRIL", "MAIO", "JUNHO",
    "JULHO", "AGOSTO", "SETEMBRO", "OUTUBRO", "NOVEMBRO", "DEZEMBRO",
]

# Índices de slide na referência (0-based) — conferidos em RMA_referencia.pptx.
SLIDE_TITULO = 0
SLIDE_OBJETIVO = 1
SLIDES_MATRIZ = (2, 9)
SLIDE_EMPRESTIMOS = 14
SLIDES_REPARAVEIS = (16, 17)
SLIDE_PAGAMENTOS = 13
SLIDE_EMPRESTIMOS_DASHBOARD = 15
SLIDE_ATRASOS_DASHBOARD = 18

VALORES_ONDE_ENCONTRA_EMPRESA = ["VEE ONE", "AMA - VEE ONE", "V1 PAMA-LS"]
VALORES_ONDE_ENCONTRA_TERCEIRIZADA = ["WILLIAM", "LEAP", "AV AERONAUTICA", "NÃO TEVE MOVIMENTAÇÃO"]
VALORES_ONDE_ENCONTRA_EMPRESA_TERCEIRIZADA = VALORES_ONDE_ENCONTRA_EMPRESA + VALORES_ONDE_ENCONTRA_TERCEIRIZADA

SLIDES_AERONAVE_INICIO = 3
SLIDES_AERONAVE_FIM = 8

# Slides da referência com conteúdo COPIADO da empresa (imagem/tabela da
# apresentação real deles, não produzido por nós) — Utilização (fonte é
# um PDF externo, fora do nosso alcance), tabela de Faturamento e Notas
# Fiscais. Removidos sempre — pedido do Wallace em 2026-07-13 ("as
# informações do slide são todas produzidas pela gente, o que tiver
# copiado tira"). Se um dia produzirmos esses dados nós mesmos, substituir
# por um slide nativo em vez de só remover.
SLIDES_COPIADOS_EXTRAS = (10, 11, 12)

COLUNAS_EMERGENCIA = [
    ("om_emg", "OM_EMG", 0.6), ("om", "OM", 0.7), ("numero_emergencia", "EMERGÊNCIA", 1.1),
    ("pn", "PN", 0.9), ("nomenclatura", "NOMENCLATURA", 1.6), ("categoria", "CAT", 0.5),
    ("matricula_aeronave", "MATR", 0.6), ("situacao", "ST_EMG", 1.0), ("tpemg", "TPEMG", 0.6),
    ("data_abertura", "DT_EMG", 0.9), ("data_info", "INFO EMG", 0.9), ("quantidade", "QT_EMG", 0.5),
    ("unidade_medida", "UE", 0.5), ("prazo_entrega", "PRAZO DE ENTREGA", 0.9),
    ("dpe", "DPE VEE ONE", 1.1), ("atendido_cancelado", "Atd/cancelada", 1.0),
    ("dias_atraso", "DIAS ATRASO", 0.6), ("dias_corridos", "DIAS CORRIDOS", 0.6),
    ("estoque", "Estoque", 0.6),
    ("retirado_empresa_recibo_obrigatorio", "Retirado pela empresa? Obrigatório Recibo.", 0.9),
    ("obs_coordenadoria_fiscal", "OBSERVAÇÃO COORDENADORIA / FISCAL", 3.0),
]
COLUNAS_EMPRESTIMO = [
    ("numero_ordem", "Nº ORDEM", 0.5), ("part_number", "PN", 0.9), ("descricao", "DESCRIÇÃO", 1.6),
    ("sn_lt", "SN/LT", 0.7), ("categoria", "CAT", 0.4), ("quantidade_texto", "QTD", 0.6),
    ("pedido_emg", "PEDIDO/EMG", 0.9), ("motivo", "MOTIVO", 1.8), ("pj", "PJ", 0.5),
    ("pedido_envio", "PEDIDO DE ENVIO", 0.8), ("anv", "ANV", 0.5), ("destino", "DESTINO", 0.7),
    ("nf_gmm", "NF/GMM", 0.7), ("rastreio", "RASTREIO", 0.7), ("pn_devolvido", "PN DEVOLVIDO", 0.8),
    ("detalhe_entrega", "STATUS (ENTREGA)", 1.3), ("numero_rc", "Nº RC", 0.6),
    ("nf_devolucao_vee_one", "NF DEVOLUÇÃO VEE ONE", 0.9), ("data_devolucao", "DATA DEVOLUÇÃO", 0.8),
    ("observacao_fiscal", "OBS. FISCAL", 1.4), ("observacao_empresa", "OBS. EMPRESA", 1.4),
    ("status", "STATUS", 0.5),
]
COLUNAS_REPARAVEL = [
    ("os", "OS", 0.6), ("pn", "PN", 0.9), ("cff", "CFF", 0.6), ("nomenclatura", "NOMENCLATURA", 1.6),
    ("sn", "SN", 0.7), ("unidade_solicitante", "UNIDADE SOLICITANTE", 0.8), ("situacao", "SITUAÇÃO", 0.8),
    ("condicao", "CONDIÇÃO", 0.9), ("onde_se_encontra", "ONDE SE ENCONTRA", 0.9),
    ("data_inicio", "DATA INÍCIO", 0.7), ("tat_siloms", "TAT SILOMS", 0.6),
    ("data_retorno_prevista", "DATA RETORNO PREVISTA", 0.8),
    ("sn_trocado_exchange", "SN TROCADO (EXCHANGE)", 0.8), ("termo_recebimento", "TERMO RECEBIMENTO", 0.9),
]
COLUNAS_PAGAMENTO = [
    ("modulo", "MÓDULO", 0.5), ("referencia", "REFERÊNCIA", 0.7),
    ("numero_nota_fiscal", "Nº NOTA FISCAL", 0.8), ("data", "DATA", 0.8),
    ("valor_nfs_fmt", "VALOR NFs", 1.1), ("faturado_fmt", "FATURADO", 1.1),
    ("pendente_fmt", "PENDENTE", 1.0), ("situacao", "SITUAÇÃO", 0.7),
    ("empenho_responsavel", "EMPENHO RESPONSÁVEL", 1.1),
]
COLUNAS_ATRASOS_ABERTAS = [
    ("numero_emergencia", "EMERGÊNCIA", 1.1), ("matricula_aeronave", "AERONAVE", 0.7),
    ("tpemg", "TIPO", 0.6), ("situacao", "SITUAÇÃO", 1.0), ("data_abertura", "ABERTURA", 0.8),
    ("prazo_entrega", "PRAZO", 0.8), ("dias_atraso", "DIAS DE ATRASO", 0.8),
]
COLUNAS_ATRASOS_ENTREGAS = [
    ("numero_emergencia", "EMERGÊNCIA", 1.1), ("matricula_aeronave", "AERONAVE", 0.7),
    ("tpemg", "TIPO", 0.6), ("data_abertura", "ABERTURA", 0.8), ("prazo_entrega", "PRAZO", 0.8),
    ("atendido_cancelado_fmt", "CANCELAMENTO/CONCLUSÃO", 1.0), ("dias_atraso", "DIAS DE ATRASO", 0.8),
    ("situacao_prazo", "SITUAÇÃO", 0.8),
]
ABREV_SEMANA = ["S", "T", "Q", "Q", "S", "S", "D"]

COR_CABECALHO = RGBColor(232, 232, 232)
COR_AMBER = RGBColor(192, 106, 0)
COR_TEXTO = RGBColor(32, 32, 32)
FONTE_TAMANHO_EMERGENCIA = Pt(6)
FONTE_TAMANHO_DIA = Pt(8)

COR_MONTADA = RGBColor(217, 234, 211)
COR_DESMONTADA = RGBColor(255, 242, 204)
COR_FORA_CONTRATO = RGBColor(255, 249, 230)
COR_BRANCO = RGBColor(255, 255, 255)
COR_FIM_DE_SEMANA_TEXTO = RGBColor(204, 0, 0)

ESPACO_ENTRE_TABELAS = Emu(500000)
FONTE_TAMANHO_MATRIZ = Pt(6)

# Cores de DADO dos gráficos (não é "brand" de slide — ver docstring acima).
COR_GRAFICO_AMBER = RGBColor(242, 169, 59)
COR_GRAFICO_CYAN = RGBColor(95, 208, 217)
COR_STATUS_BOM = RGBColor(79, 180, 119)
COR_STATUS_CRITICO = RGBColor(226, 86, 79)

# Fonte/cor usadas na renderização em imagem da matriz (Tabela 1.2) — ver
# _renderizar_matriz_imagem.
FONTE_TABELA_IMAGEM = "/System/Library/Fonts/Supplemental/Arial.ttf"
COR_CABECALHO_IMAGEM = (232, 232, 232)


def _formatar_valor(valor):
    if valor is None or (isinstance(valor, (float, pd.Timestamp, type(pd.NaT))) and pd.isna(valor)):
        return ""
    if isinstance(valor, (pd.Timestamp, datetime, date)):
        return valor.strftime("%d/%m/%Y")
    if isinstance(valor, float) and valor.is_integer():
        return str(int(valor))
    return str(valor)


def _carregar_dados_emprestimos(ano, mes):
    """Empréstimos (planilha "Devoluções") — devolve (do_mes, todos). Filtra
    `do_mes` por `pedido_envio` (data em que o item saiu do estoque),
    pedido do Wallace em 2026-07-12 ("os emprestimos a gente coloca os do
    mes de referencia"); `todos` é a base inteira sem filtro de mês, pra
    "uma estatística total" (pedido do Wallace, mesma conversa). Ordenado
    por numero_ordem (ordem cronológica da planilha). Cada linha vira
    "quantidade_efetiva" pra estatística ponderada pela quantidade real
    (pedido do Wallace: "a quantidade multiplica la pela quantidade ...
    as vezes uma linha tem 10 ea") — linha sem quantidade registrada conta
    como 1."""
    df = pd.read_excel(DADOS_TRATADOS / "base_devolucoes_tratada.xlsx")
    df["pedido_envio_dt"] = pd.to_datetime(df["pedido_envio"], errors="coerce")
    df["quantidade_efetiva"] = df["quantidade"].fillna(1)
    df = df.sort_values("numero_ordem")
    do_mes = df[(df["pedido_envio_dt"].dt.year == ano) & (df["pedido_envio_dt"].dt.month == mes)].copy()
    return do_mes, df


def _carregar_dados_reparaveis():
    """Reparáveis — devolve (todas_em_aberto, com_empresa_ou_terceirizada).
    "Em aberto" é o filtro padrão da própria tela do site (reparaveis.py);
    o 2º filtro (onde_se_encontra) é o pedido do Wallace em 2026-07-12,
    "primeiro todas e depois filtra colocando os que ta com a empresa ou
    nas tercerizadas, igual os filtros do slide [...] onde se encontra" —
    ver VALORES_ONDE_ENCONTRA_EMPRESA_TERCEIRIZADA."""
    df = pd.read_excel(DADOS_TRATADOS / "base_reparaveis_tratada.xlsx")
    abertas = df[df["em_aberto"]].copy()
    filtradas = abertas[abertas["onde_se_encontra"].isin(VALORES_ONDE_ENCONTRA_EMPRESA_TERCEIRIZADA)].copy()
    return abertas, filtradas


def _carregar_dados_pagamentos():
    """Pagamentos — mesma fonte e mesmas colunas exibidas da tela do site
    (pagamentos.py) — pedido do Wallace em 2026-07-12 ("os pagamentos vc
    pega no nosso drive"). `base_pagamentos_tratada.xlsx` tem 3 abas —
    "Pagamentos" (1 linha por nota fiscal) e "Contrato" (resumo, 1 linha)."""
    caminho = DADOS_TRATADOS / "base_pagamentos_tratada.xlsx"
    df = pd.read_excel(caminho, sheet_name="Pagamentos")
    contrato = pd.read_excel(caminho, sheet_name="Contrato").iloc[0]
    _fmt = lambda v: f"R$ {v:,.2f}" if pd.notna(v) else ""
    df["valor_nfs_fmt"] = df["valor_nfs"].apply(_fmt)
    df["faturado_fmt"] = df["faturado"].apply(_fmt)
    df["pendente_fmt"] = df["pendente"].apply(_fmt)
    return df, contrato


def _carregar_dados_atrasos(ano, mes):
    """Atrasos — mesma lógica de fechamento_mensal.py `_atrasos()` (regras
    do Wallace, 2026-07-10, ver 00_Instrucoes/atrasos.md), pedido pra virar
    nativo aqui em 2026-07-12 ("trabalhar os atrasos tb, igual do site").
    Devolve (abertas, concluidas_mes, resumo)."""
    df = pd.read_excel(DADOS_TRATADOS / "historico_completo_emergencias.xlsx").copy()
    df["atendido_cancelado_dt"] = pd.to_datetime(df["atendido_cancelado"], errors="coerce")

    abertas = df[df["em_aberto"]].sort_values("dias_atraso", ascending=False).copy()

    inicio_mes = pd.Timestamp(ano, mes, 1)
    fim_mes = inicio_mes + pd.offsets.MonthEnd(0)
    concluidas_mes = df[
        (~df["em_aberto"])
        & (df["atendido_cancelado_dt"] >= inicio_mes)
        & (df["atendido_cancelado_dt"] <= fim_mes)
    ].copy()
    concluidas_mes["situacao_prazo"] = concluidas_mes["dias_atraso"].apply(lambda d: "No prazo" if d <= 0 else "Atrasado")
    concluidas_mes["atendido_cancelado_fmt"] = concluidas_mes["atendido_cancelado_dt"].dt.strftime("%d/%m/%Y")

    total_previstas = len(concluidas_mes)
    no_prazo = int((concluidas_mes["dias_atraso"] <= 0).sum())
    resumo = {
        "periodo": f"{inicio_mes.strftime('%d/%m')} - {fim_mes.strftime('%d/%m')}",
        "total_previstas": total_previstas,
        "no_prazo": no_prazo,
        "atrasadas": total_previstas - no_prazo,
        "pct": (100 * no_prazo / total_previstas) if total_previstas else 0.0,
    }
    return abertas, concluidas_mes, resumo


def _carregar_dados_mes(ano, mes):
    """Junta Cômputo Mensal (matriz + motivos) com o histórico completo de
    emergências — devolve um dict {matricula: {"emergencias": DataFrame,
    "dia_anv": Series}}, na ordem em que aparecem no motivos.csv (mesma
    ordem cronológica da negativação)."""
    pasta = DADOS_TRATADOS / "computo_mensal"
    mes_ref = f"{ano}-{mes:02d}"
    df_matriz = pd.read_csv(pasta / f"{mes_ref}_matriz.csv", dtype={"matricula": str})
    df_motivos = pd.read_csv(pasta / f"{mes_ref}_motivos.csv", dtype={"matricula": str})
    with open(pasta / f"{mes_ref}_resumo.json", encoding="utf-8") as f:
        resumo = json.load(f)

    df_hist = pd.read_excel(DADOS_TRATADOS / "historico_completo_emergencias.xlsx", dtype={"matricula_aeronave": str})

    aeronaves = {}
    for matricula in df_motivos["matricula"].drop_duplicates():
        numeros = df_motivos.loc[df_motivos["matricula"] == matricula, "numero_emergencia"].unique().tolist()
        emergencias = df_hist[df_hist["numero_emergencia"].isin(numeros)]
        dia_anv = df_matriz.loc[df_matriz["matricula"] == matricula].sort_values("dia")
        aeronaves[matricula] = {"emergencias": emergencias, "dia_anv": dia_anv}
    return aeronaves, resumo


def _remover_slides(prs, indice_inicio, indice_fim):
    """Remove os slides de índice_inicio até índice_fim (inclusive), na
    ordem correta (de trás pra frente, pra não invalidar os índices). Tira
    também a relação (rId) do slide, não só da lista de ordem — senão o
    part continua "alcançável" no pacote e o arquivo salva com partes
    duplicadas (slide novo reaproveitando um nome de arquivo ainda em uso)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for indice in range(indice_fim, indice_inicio - 1, -1):
        elemento = slides[indice]
        prs.part.drop_rel(elemento.rId)
        xml_slides.remove(elemento)


def _slide_em_branco(prs, indice_apos):
    """Insere um slide em branco logo depois do índice dado, devolve o
    objeto slide. Usa o layout do 1º slide da apresentação (mantém a mesma
    "master" visual).

    IMPORTANTE — ordem de chamada: o python-pptx nomeia o arquivo XML do
    slide novo só com base na CONTAGEM atual de slides (`slide%d.xml`, sem
    checar se esse nome já existe de verdade no pacote). Se a gente já
    tivesse apagado os slides antigos antes de chamar isso, a contagem cai
    e o nome novo colide com um slide intocado mais adiante (ex.: slide14,
    slide15...), corrompendo o arquivo salvo ("Duplicate name"). Por isso
    `gerar_apresentacao()` sempre ADICIONA todos os slides novos primeiro
    (com a contagem ainda "cheia", nomes novos só no fim, sem colisão) e só
    depois apaga os antigos."""
    layout = prs.slides[0].slide_layout
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    xml_slides = prs.slides._sldIdLst
    novo_elemento = xml_slides[-1]
    xml_slides.remove(novo_elemento)
    xml_slides.insert(indice_apos + 1, novo_elemento)
    return slide


COR_BORDA = "808080"


def _fechar_bordas(celula, cor_hex=COR_BORDA, espessura_pt=0.5):
    """Desenha as 4 bordas da célula explicitamente (lnL/lnR/lnT/lnB) —
    sem isso, dependendo do visualizador, célula com fundo customizado pode
    aparecer sem grade ("aberta"), pedido do Wallace ("tudo fechado").

    IMPORTANTE — ordem do XML: o schema de CT_TableCellProperties exige
    lnL/lnR/lnT/lnB ANTES do preenchimento (fill) da célula, senão o
    PowerPoint recusa/"repara" o arquivo. Como o fundo (`cor_fundo`) já foi
    definido antes desta função ser chamada, os 4 elementos são inseridos
    no INÍCIO de tcPr (índices 0-3), empurrando o fill (se houver) pro
    final — nunca com `append`, que os colocaria depois do fill."""
    tcPr = celula._tc.get_or_add_tcPr()
    espessura_emu = int(Pt(espessura_pt))
    for i, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
        existente = tcPr.find(qn(tag))
        if existente is not None:
            tcPr.remove(existente)
        ln = tcPr.makeelement(qn(tag), {
            "w": str(espessura_emu), "cap": "flat", "cmpd": "sng", "algn": "ctr",
        })
        solid = ln.makeelement(qn("a:solidFill"), {})
        srgb = solid.makeelement(qn("a:srgbClr"), {"val": cor_hex})
        solid.append(srgb)
        ln.append(solid)
        tcPr.insert(i, ln)


def _preencher_celula(celula, texto, negrito=False, tamanho=Pt(8), cor_fundo=None,
                       cor_texto=None, bordas=True, margem_celula=None):
    celula.text = texto
    celula.text_frame.word_wrap = True
    celula.vertical_anchor = MSO_ANCHOR.MIDDLE
    if margem_celula is not None:
        celula.margin_left = margem_celula
        celula.margin_right = margem_celula
        celula.margin_top = margem_celula
        celula.margin_bottom = margem_celula
    if cor_fundo is not None:
        celula.fill.solid()
        celula.fill.fore_color.rgb = cor_fundo
    for paragrafo in celula.text_frame.paragraphs:
        paragrafo.alignment = PP_ALIGN.CENTER
        for run in paragrafo.runs:
            run.font.size = tamanho
            run.font.bold = negrito
            run.font.color.rgb = cor_texto or COR_TEXTO
    if bordas:
        _fechar_bordas(celula)


def _montar_slide_aeronave(slide, matricula, dados_aeronave, largura_total, ano, mes):
    emergencias = dados_aeronave["emergencias"]
    dia_anv = dados_aeronave["dia_anv"]

    linhas_dia = dia_anv.to_dict("records")
    n_dias = len(linhas_dia)

    topo = Emu(150000)
    tabela_dia_shape = slide.shapes.add_table(2, n_dias + 1, Emu(300000), topo, Emu(largura_total - 600000), Emu(250000))
    tabela_dia = tabela_dia_shape.table
    _preencher_celula(tabela_dia.cell(0, 0), "DIA\nANV", negrito=True, tamanho=FONTE_TAMANHO_DIA, cor_fundo=COR_CABECALHO)
    _preencher_celula(tabela_dia.cell(1, 0), str(matricula), negrito=True, tamanho=FONTE_TAMANHO_DIA, cor_fundo=COR_CABECALHO)
    for i, linha in enumerate(linhas_dia):
        dia = int(linha["dia"])
        semana = ABREV_SEMANA[date(ano, mes, dia).weekday()]
        _preencher_celula(tabela_dia.cell(0, i + 1), f"{dia}\n{semana}", tamanho=FONTE_TAMANHO_DIA)
        valor = linha.get("montada")
        if pd.isna(valor):
            texto, cor = "", None
        else:
            texto, cor = ("1", COR_MONTADA) if valor == 1 else ("0", COR_DESMONTADA)
        _preencher_celula(tabela_dia.cell(1, i + 1), texto, cor_fundo=cor, tamanho=FONTE_TAMANHO_DIA)

    topo_emergencias = topo + Emu(250000) + ESPACO_ENTRE_TABELAS
    colunas = COLUNAS_EMERGENCIA
    n_linhas = len(emergencias) + 1
    largura_relativa_total = sum(c[2] for c in colunas)
    tabela_emg_shape = slide.shapes.add_table(
        n_linhas, len(colunas), Emu(300000), topo_emergencias,
        Emu(largura_total - 600000), Emu(220000 * n_linhas),
    )
    tabela_emg = tabela_emg_shape.table
    for i, (campo, rotulo, peso) in enumerate(colunas):
        tabela_emg.columns[i].width = Emu(int(largura_relativa_total and (largura_total - 600000) * peso / largura_relativa_total))
        _preencher_celula(tabela_emg.cell(0, i), rotulo, negrito=True, tamanho=FONTE_TAMANHO_EMERGENCIA, cor_fundo=COR_CABECALHO)
    for i, (_, registro) in enumerate(emergencias.iterrows(), start=1):
        for j, (campo, _, _) in enumerate(colunas):
            _preencher_celula(tabela_emg.cell(i, j), _formatar_valor(registro.get(campo)), tamanho=FONTE_TAMANHO_EMERGENCIA)


def _atualizar_titulo_objetivo(prs, mes):
    """Slide 1 (título "RMA {MÊS}") e slide 2 (objetivo, "...de {mês}") — só
    troca o nome do mês, preservando tudo o mais (fonte, tamanho, o resto do
    texto) exatamente como na referência. O mês já é um RUN isolado nos dois
    casos (conferido em RMA_referencia.pptx), então dá pra trocar o texto do
    run direto sem tocar em mais nada."""
    mes_maiusculo = MESES_PT_MAIUSCULO[mes - 1]
    mes_minusculo = mes_maiusculo.lower()
    nomes_meses_maiusculo = set(MESES_PT_MAIUSCULO)
    nomes_meses_minusculo = {m.lower() for m in MESES_PT_MAIUSCULO}
    for indice_slide, texto_novo, condicao in (
        (SLIDE_TITULO, mes_maiusculo, lambda t: t.strip() in nomes_meses_maiusculo),
        (SLIDE_OBJETIVO, mes_minusculo, lambda t: t.strip() in nomes_meses_minusculo),
    ):
        slide = prs.slides[indice_slide]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragrafo in shape.text_frame.paragraphs:
                for run in paragrafo.runs:
                    if condicao(run.text):
                        run.text = texto_novo


def _remover_imagens(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            shape._element.getparent().remove(shape._element)


def _renderizar_matriz_imagem(df_matriz_mes, pontuadas, fora_listadas, resumo, ano, mes, caminho_png,
                               largura_px=2400, largura_rotulo_px=90, altura_linha_px=26, altura_cabecalho_px=40):
    """Desenha a Tabela 1.2 (matriz aeronave x dia) como IMAGEM, produzida
    por nós a partir do Cômputo Mensal (`resumo`/`df_matriz_mes`, mesmos
    dados da aba 1.2 da RMA oficial) — pedido do Wallace em 2026-07-13
    ("a imagem [...] ficou ruim, vamos colocar como uma foto que a gente
    produz com base no fechamento e com base na aba 1.2 da rma"). Antes
    disso era uma tabela nativa do pptx (ficou feia com 30 colunas
    espremidas); como imagem dá pra controlar o layout com precisão,
    igual as tabelas de anexo da Ata de Reunião (Pillow, sem lib pesada)."""
    try:
        fonte = ImageFont.truetype(FONTE_TABELA_IMAGEM, 12)
        fonte_negrito = ImageFont.truetype(FONTE_TABELA_IMAGEM.replace("Arial.ttf", "Arial Bold.ttf"), 12)
    except OSError:
        fonte = fonte_negrito = ImageFont.load_default()

    ultimo_dia = resumo["ultimo_dia_mes"]
    todas_aeronaves = list(pontuadas) + list(fora_listadas)
    largura_dia_px = (largura_px - largura_rotulo_px) // ultimo_dia
    largura_px = largura_rotulo_px + largura_dia_px * ultimo_dia
    n_linhas_dado = len(todas_aeronaves) + 1  # aeronaves + "Média Diária"
    altura_px = altura_cabecalho_px + n_linhas_dado * altura_linha_px

    img = Image.new("RGB", (largura_px, altura_px), "white")
    desenho = ImageDraw.Draw(img)

    def _celula(x, y, w, h, cor_fundo, texto=None, fnt=None, cor_texto=(20, 20, 20), centralizado=True):
        if cor_fundo:
            desenho.rectangle([x, y, x + w, y + h], fill=cor_fundo)
        desenho.rectangle([x, y, x + w, y + h], outline=(190, 190, 190))
        if texto:
            fnt = fnt or fonte
            bbox = desenho.textbbox((0, 0), texto, font=fnt)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            tx = x + (w - tw) // 2 if centralizado else x + 4
            ty = y + (h - th) // 2 - bbox[1]
            desenho.text((tx, ty), texto, font=fnt, fill=cor_texto)

    _celula(0, 0, largura_rotulo_px, altura_cabecalho_px, COR_CABECALHO_IMAGEM, "DIA/ANV", fonte_negrito)
    for dia in range(1, ultimo_dia + 1):
        wd = date(ano, mes, dia).weekday()
        cor_texto = (204, 0, 0) if wd >= 5 else (20, 20, 20)
        x = largura_rotulo_px + (dia - 1) * largura_dia_px
        _celula(x, 0, largura_dia_px, altura_cabecalho_px, COR_CABECALHO_IMAGEM, f"{dia}\n{ABREV_SEMANA[wd]}", fonte_negrito, cor_texto)

    pivot = df_matriz_mes.set_index(["matricula", "dia"])["montada"] if not df_matriz_mes.empty else None
    y = altura_cabecalho_px
    for matricula in todas_aeronaves:
        _celula(0, y, largura_rotulo_px, altura_linha_px, COR_CABECALHO_IMAGEM, str(matricula), fonte_negrito)
        for dia in range(1, ultimo_dia + 1):
            x = largura_rotulo_px + (dia - 1) * largura_dia_px
            if matricula in fora_listadas:
                _celula(x, y, largura_dia_px, altura_linha_px, (255, 249, 230))
                continue
            valor = pivot.get((matricula, dia)) if pivot is not None else None
            if valor is None or pd.isna(valor):
                _celula(x, y, largura_dia_px, altura_linha_px, None)
            else:
                cor = (217, 234, 211) if valor == 1 else (255, 242, 204)
                _celula(x, y, largura_dia_px, altura_linha_px, cor, str(int(valor)), fonte)
        y += altura_linha_px

    # "Média Diária" calculada por nós, direto da matriz (não depende de
    # `calcular_computo_mensal.py` ter salvo um resumo pré-calculado — esse
    # script parou de gravar esses campos extras, ver docstring do módulo).
    medias_diarias = (
        df_matriz_mes.dropna(subset=["montada"]).groupby("dia")["montada"].mean() * 100
        if not df_matriz_mes.empty else pd.Series(dtype=float)
    )
    _celula(0, y, largura_rotulo_px, altura_linha_px, COR_CABECALHO_IMAGEM, "Média Diária", fonte_negrito)
    for dia in range(1, ultimo_dia + 1):
        x = largura_rotulo_px + (dia - 1) * largura_dia_px
        valor = medias_diarias.get(dia)
        texto = f"{valor:.0f}%" if valor is not None and not pd.isna(valor) else "—"
        _celula(x, y, largura_dia_px, altura_linha_px, COR_CABECALHO_IMAGEM, texto, fonte)

    caminho_png.parent.mkdir(parents=True, exist_ok=True)
    img.save(caminho_png)
    return caminho_png, largura_px, altura_px


def _construir_tabela_matriz(slide, df_matriz_mes, pontuadas, fora_listadas, resumo, indicadores_oficiais, ano, mes, largura_total, altura_total, caminho_temp):
    """Insere a Tabela 1.2 (matriz aeronave x dia + MMAM/P/PMAX/IFD) como
    imagem produzida por nós (ver `_renderizar_matriz_imagem`), no lugar
    da imagem copiada da referência. Os 4 números do resumo (MMAM/P/PMAX/
    IFD) vêm de `indicadores_oficiais` — a mesma planilha "RMA em
    andamento" oficial usada pela Ata de Reunião (`gerar_ata_reuniao.py`),
    não de `calcular_computo_mensal.py` (que parou de salvar esses campos
    e, de qualquer forma, tem uma pequena divergência conhecida vs. o
    número oficial da empresa — ver computo_mensal.md)."""
    _remover_imagens(slide)

    largura_imagem_emu = largura_total - 240000
    caminho_png, largura_px, altura_px = _renderizar_matriz_imagem(df_matriz_mes, pontuadas, fora_listadas, resumo, ano, mes, caminho_temp)
    altura_imagem_emu = int(largura_imagem_emu * altura_px / largura_px)
    slide.shapes.add_picture(str(caminho_png), Emu(120000), Emu(1300000), width=Emu(largura_imagem_emu))
    caminho_png.unlink()

    RESERVA_RESUMO = 750000
    if altura_imagem_emu + RESERVA_RESUMO > int(altura_total):
        altura_imagem_emu = int(altura_total) - RESERVA_RESUMO

    caixa_resumo = slide.shapes.add_textbox(Emu(120000), Emu(1300000 + altura_imagem_emu + 150000), Emu(largura_total - 240000), Emu(600000))
    linhas_resumo = [
        ("Média Mensal de Aeronaves Montadas MMAM", f"{indicadores_oficiais['mmam'] * 100:.2f}%".replace(".", ",")),
        ("PONTUAÇÃO OBTIDA (P)", f"{indicadores_oficiais['pontuacao_obtida']:.0f}"),
        ("PONT. MAX NO MÊS (PMAX)", f"{indicadores_oficiais['pont_max']:.0f}"),
        ("Índice Final de Desempenho IFD", f"{indicadores_oficiais['ifd']:.2f}".replace(".", ",")),
    ]
    tf = caixa_resumo.text_frame
    for i, (rotulo, valor) in enumerate(linhas_resumo):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{rotulo}: {valor}"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = COR_AMBER


def _titulo_slide(slide, texto, largura_total, margem=Emu(340000), topo=Emu(340000), tamanho=Pt(20)):
    """Caixa de texto simples com um título em negrito — usada acima das
    tabelas de Empréstimos/Reparáveis. Devolve o topo livre logo abaixo."""
    caixa = slide.shapes.add_textbox(margem, topo, Emu(largura_total) - 2 * margem, Emu(380000))
    p = caixa.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = tamanho
    run.font.bold = True
    run.font.color.rgb = COR_TEXTO
    return topo + Emu(380000)


def _construir_tabela_dados(slide, df, colunas, largura_total, topo, margem, tamanho_fonte, altura_linha, margem_celula=None):
    """Tabela genérica (cabeçalho + 1 linha por registro), mesmo padrão da
    tabela de emergências dos slides por aeronave — reaproveitada aqui pra
    Empréstimos e Reparáveis. `margem_celula` (opcional) zera o preenchimento
    interno da célula — usado no modo "densidade máxima"."""
    n_linhas = len(df) + 1
    largura_relativa_total = sum(c[2] for c in colunas)
    tabela_shape = slide.shapes.add_table(
        n_linhas, len(colunas), margem, topo, Emu(largura_total) - 2 * margem, Emu(altura_linha) * n_linhas,
    )
    tabela = tabela_shape.table
    for i, (campo, rotulo, peso) in enumerate(colunas):
        tabela.columns[i].width = Emu(int((largura_total - 2 * margem) * peso / largura_relativa_total))
        _preencher_celula(tabela.cell(0, i), rotulo, negrito=True, tamanho=tamanho_fonte, cor_fundo=COR_CABECALHO, margem_celula=margem_celula)
    for i, (_, registro) in enumerate(df.iterrows(), start=1):
        for j, (campo, _, _) in enumerate(colunas):
            _preencher_celula(tabela.cell(i, j), _formatar_valor(registro.get(campo)), tamanho=tamanho_fonte, margem_celula=margem_celula)
    return tabela_shape


def _legenda_grafico(slide, texto, largura, left, top):
    caixa = slide.shapes.add_textbox(left, top, largura, Emu(280000))
    p = caixa.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COR_TEXTO


def _grafico_barra(slide, categorias, valores, cor, left, top, largura, altura, horizontal=False):
    """Gráfico de barras simples (1 série, 1 cor) — mesmo estilo dos
    gráficos "Por categoria"/"Distribuição por condição" do site."""
    dados_grafico = CategoryChartData()
    dados_grafico.categories = [str(c) for c in categorias]
    dados_grafico.add_series("Quantidade", valores)
    tipo = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    grafico_frame = slide.shapes.add_chart(tipo, left, top, largura, altura, dados_grafico)
    grafico = grafico_frame.chart
    grafico.has_legend = False
    serie = grafico.plots[0].series[0]
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = cor
    grafico.plots[0].has_data_labels = True
    grafico.plots[0].data_labels.font.size = Pt(9)
    grafico.category_axis.tick_labels.font.size = Pt(9)
    grafico.value_axis.has_major_gridlines = True
    grafico.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE3, 0xE3, 0xE3)
    grafico.value_axis.major_gridlines.format.line.width = Pt(0.5)
    return grafico_frame


def _grafico_rosca(slide, categorias, valores, cores, left, top, largura, altura):
    """Gráfico de rosca (donut) com 1 cor por fatia — igual ao gráfico de
    Status (Pendente/OK) do site."""
    dados_grafico = CategoryChartData()
    dados_grafico.categories = [str(c) for c in categorias]
    dados_grafico.add_series("Quantidade", valores)
    grafico_frame = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, largura, altura, dados_grafico)
    grafico = grafico_frame.chart
    grafico.has_legend = True
    grafico.legend.position = XL_LEGEND_POSITION.BOTTOM
    grafico.legend.include_in_layout = False
    grafico.legend.font.size = Pt(9)
    grafico.plots[0].has_data_labels = True
    for ponto, cor in zip(grafico.plots[0].series[0].points, cores):
        ponto.format.fill.solid()
        ponto.format.fill.fore_color.rgb = cor
    return grafico_frame


def _construir_slide_emprestimos_estatisticas(slide, df, titulo, largura_total):
    """Estatísticas de Empréstimos — mesmos 3 gráficos de emprestimos.py
    (Status, Por categoria, Top destinos), reaproveitada tanto pro mês de
    referência quanto pro total (pedido do Wallace, 2026-07-12: "vamos
    mostrar estatística do mês, a lista completa do mês e uma estatística
    total"). Pesado pela QUANTIDADE de cada linha, não só a contagem de
    linhas (pedido do Wallace, mesma conversa: "a quantidade multiplica la
    pela quantidade ... as vezes uma linha tem 10 ea") — mesmo ajuste
    espelhado em emprestimos.py no site."""
    topo = _titulo_slide(slide, titulo, largura_total, topo=Emu(150000), tamanho=Pt(20))

    total = len(df)
    total_qtd = df["quantidade_efetiva"].sum()
    pendentes = int((df["status"] == "Pendente").sum())
    ok = int((df["status"] == "OK").sum())
    caixa = slide.shapes.add_textbox(Emu(340000), topo + Emu(120000), Emu(largura_total - 680000), Emu(400000))
    p = caixa.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"{total} itens (linhas)  ·  {total_qtd:,.0f}".replace(",", ".") +
        f" de quantidade  ·  {pendentes} pendentes  ·  {ok} OK"
    )
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COR_AMBER

    largura_grafico = Emu((largura_total - 4 * 200000) // 3)
    top_grafico = topo + Emu(460000)

    _legenda_grafico(slide, "Status (por quantidade)", largura_grafico, Emu(200000), top_grafico)
    status_contagem = df.groupby("status")["quantidade_efetiva"].sum()
    cores_status = [COR_STATUS_CRITICO if s == "Pendente" else COR_STATUS_BOM for s in status_contagem.index]
    _grafico_rosca(slide, status_contagem.index, status_contagem.values, cores_status,
                   Emu(200000), top_grafico + Emu(320000), largura_grafico, Emu(3200000))

    left2 = Emu(200000) + largura_grafico + Emu(200000)
    _legenda_grafico(slide, "Por categoria (C/T/R) — quantidade", largura_grafico, left2, top_grafico)
    cat_contagem = df.groupby("categoria")["quantidade_efetiva"].sum()
    _grafico_barra(slide, cat_contagem.index, cat_contagem.values, COR_GRAFICO_AMBER,
                   left2, top_grafico + Emu(320000), largura_grafico, Emu(3200000))

    left3 = left2 + largura_grafico + Emu(200000)
    _legenda_grafico(slide, "Top destinos — quantidade", largura_grafico, left3, top_grafico)
    destino_contagem = df.groupby("destino")["quantidade_efetiva"].sum().sort_values(ascending=False).head(8)
    _grafico_barra(slide, destino_contagem.index, destino_contagem.values, COR_GRAFICO_CYAN,
                   left3, top_grafico + Emu(320000), largura_grafico, Emu(3200000), horizontal=True)


def _construir_slide_atrasos_situacao(slide, abertas, largura_total):
    """Slide "Atrasos — Situação Atual" — mesmas 3 métricas + tabela da
    seção "Situação atual (em aberto agora)" do site (fechamento_mensal.py,
    _atrasos). Cabem as 14 linhas inteiras (não precisa de amostra)."""
    topo = _titulo_slide(slide, "Atrasos — Situação Atual (em aberto agora)", largura_total, topo=Emu(150000), tamanho=Pt(20))
    total = len(abertas)
    atrasadas = int((abertas["dias_atraso"] > 0).sum())
    no_prazo = total - atrasadas

    caixa = slide.shapes.add_textbox(Emu(340000), topo + Emu(120000), Emu(largura_total - 680000), Emu(400000))
    run = caixa.text_frame.paragraphs[0].add_run()
    run.text = f"{total} em aberto (VEE ONE)  ·  {no_prazo} dentro do prazo  ·  {atrasadas} atrasadas"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COR_AMBER

    _construir_tabela_dados(slide, abertas, COLUNAS_ATRASOS_ABERTAS, largura_total,
                             topo + Emu(460000), Emu(300000), Pt(10), 300000)


def _construir_slide_atrasos_entregas(slide, concluidas_mes, resumo, largura_total, n_amostra=10):
    """Slide "Atrasos — Entregas no mês" — resumo + rosca (No prazo/
    Atrasado) + amostra, mesma seção "Entregas no mês de referência" do
    site. Formato do resumo pedido pelo Wallace em 2026-07-10 (ver
    00_Instrucoes/atrasos.md). Alturas fixas e conferidas pra caber tudo
    (ver bug parecido — gráfico sobrepondo tabela — corrigido nos slides de
    Reparáveis, 2026-07-12)."""
    ALTURA_RESUMO = Emu(300000)
    ALTURA_LEGENDA = Emu(320000)
    ALTURA_GRAFICO = Emu(2200000)
    ALTURA_LINHA_TABELA = Emu(220000)
    GAP = Emu(150000)

    topo = _titulo_slide(slide, "Atrasos — Entregas no Mês de Referência", largura_total, topo=Emu(150000), tamanho=Pt(20))

    pct_fmt = f"{resumo['pct']:.2f}".replace(".", ",") + "%"
    resumo_df = pd.DataFrame([{
        "Período de apuração": resumo["periodo"],
        "Total de entregas previstas": resumo["total_previstas"],
        "Entregas no Prazo": resumo["no_prazo"],
        "QTD Mensal (%)": pct_fmt,
    }])
    _construir_tabela_dados(
        slide, resumo_df,
        [(c, c, 1.0) for c in resumo_df.columns],
        largura_total, topo + GAP, Emu(300000), Pt(10), int(ALTURA_RESUMO),
    )

    top_grafico = topo + GAP + ALTURA_RESUMO + GAP
    _legenda_grafico(slide, "Situação das entregas do mês", Emu(3200000), Emu(300000), top_grafico)
    _grafico_rosca(
        slide, ["No prazo", "Atrasado"], [resumo["no_prazo"], resumo["atrasadas"]],
        [COR_STATUS_BOM, COR_STATUS_CRITICO],
        Emu(300000), top_grafico + ALTURA_LEGENDA, Emu(3200000), ALTURA_GRAFICO,
    )

    top_tabela = top_grafico
    amostra = concluidas_mes.head(n_amostra)
    _construir_tabela_dados(
        slide, amostra, COLUNAS_ATRASOS_ENTREGAS, largura_total,
        top_tabela, Emu(3200000) + Emu(400000), Pt(9), int(ALTURA_LINHA_TABELA),
    )


def _construir_slide_pagamentos(slide, df, contrato, largura_total):
    """Slide de Pagamentos — pedido do Wallace em 2026-07-12 ("os pagamentos
    vc pega no nosso drive"): sai a imagem da referência, entra a tabela
    nativa (29 notas fiscais, cabem numa lâmina só) + o resumo do contrato."""
    _remover_imagens(slide)
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    topo = _titulo_slide(slide, "Pagamentos", largura_total, topo=Emu(150000), tamanho=Pt(20))
    caixa = slide.shapes.add_textbox(Emu(340000), topo + Emu(120000), Emu(largura_total - 680000), Emu(400000))
    run = caixa.text_frame.paragraphs[0].add_run()
    run.text = (
        f"Valor total do contrato: R$ {contrato['valor_total_contrato']:,.2f}   ·   "
        f"Liquidado: R$ {contrato['valor_liquidado']:,.2f}   ·   "
        f"Saldo a faturar: R$ {contrato['saldo_a_faturar']:,.2f}"
    )
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = COR_AMBER

    _construir_tabela_dados(slide, df, COLUNAS_PAGAMENTO, largura_total, topo + Emu(460000), Emu(300000), Pt(9), 190000)


def _paginar(df, linhas_primeira, linhas_seguintes):
    """Quebra `df` em quantas páginas forem necessárias — 1ª página com
    `linhas_primeira` linhas, as seguintes com `linhas_seguintes` (a 1ª
    tem menos espaço porque reserva lugar pro número em destaque). Sempre
    devolve pelo menos 1 página (mesmo vazia)."""
    paginas = []
    inicio = 0
    primeira = True
    total = len(df)
    while True:
        tamanho = linhas_primeira if primeira else linhas_seguintes
        fim = min(inicio + max(tamanho, 1), total)
        paginas.append(df.iloc[inicio:fim])
        primeira = False
        inicio = fim
        if inicio >= total:
            break
    if not paginas:
        paginas = [df.iloc[0:0]]
    return paginas


def _construir_paginas_tabela(prs, indice_apos, df, colunas, titulo_base, largura_total, tamanho_fonte, altura_linha,
                               rotulo_destaque=None, margem_celula=None, tamanho_destaque=Pt(28), altura_caixa_destaque=Emu(560000)):
    """Insere quantas lâminas forem necessárias pra mostrar TODAS as linhas
    de `df` (pedido do Wallace em 2026-07-12: "vamos colocar todas OS e
    todas empréstimo do mês, se não couber a gente divide o slide") — a 1ª
    página tem o número total em destaque (se `rotulo_destaque` for
    passado, ex. "OS" — igual ao cartão do site), as seguintes só a tabela
    continuando, com "(pág. N/M)" no título. Devolve o índice da última
    lâmina inserida (útil pra quem for inserir mais coisa logo depois)."""
    margem = Emu(150000)
    altura_disponivel = int(Emu(6358000))  # slide 7,5" (altura padrão 4:3) menos margens de topo/rodapé
    topo_base = Emu(150000) + Emu(460000)
    topo_destaque = topo_base + altura_caixa_destaque

    # `_construir_tabela_dados` sempre soma +1 linha de cabeçalho por conta
    # própria, então o orçamento de linhas de DADO aqui já desconta essa
    # linha extra.
    linhas_primeira = max(int((altura_disponivel - int(topo_destaque)) / altura_linha) - 1, 1)
    linhas_seguintes = max(int((altura_disponivel - int(topo_base)) / altura_linha) - 1, 1)

    paginas = _paginar(df, linhas_primeira, linhas_seguintes)
    n_paginas = len(paginas)

    indice = indice_apos
    for i, pagina_df in enumerate(paginas):
        slide = _slide_em_branco(prs, indice)
        indice += 1
        sufixo = f" (pág. {i + 1}/{n_paginas})" if n_paginas > 1 else ""
        topo = _titulo_slide(slide, titulo_base + sufixo, largura_total, topo=Emu(150000), tamanho=Pt(20))
        if i == 0 and rotulo_destaque:
            caixa = slide.shapes.add_textbox(Emu(340000), topo + Emu(120000), Emu(largura_total - 680000), altura_caixa_destaque)
            run = caixa.text_frame.paragraphs[0].add_run()
            run.text = f"{len(df)} {rotulo_destaque}"
            run.font.size = tamanho_destaque
            run.font.bold = True
            run.font.color.rgb = COR_AMBER
            topo_tabela = topo_destaque
        else:
            topo_tabela = topo_base
        _construir_tabela_dados(slide, pagina_df, colunas, largura_total, topo_tabela, margem, tamanho_fonte, altura_linha, margem_celula=margem_celula)

    return indice


def _construir_slide_reparaveis_estatisticas(slide, df_abertas, largura_total):
    """Reparáveis — só estatística agora, sem lista completa (pedido do
    Wallace em 2026-07-12: "só estatísticas, total, quantidade com ele ou
    terceirizados ... não precisa mandar a lista completa dos reparáveis
    não"). Números em destaque (total/com a empresa/terceirizadas/nas
    bases FAB, igual ao estilo de cartão do site) + 2 gráficos: "Onde se
    encontra" (as 3 categorias) e "Distribuição por condição"
    (reaproveita reparaveis.py)."""
    topo = _titulo_slide(slide, "Reparáveis — Estatísticas", largura_total, topo=Emu(150000), tamanho=Pt(20))

    total = len(df_abertas)
    n_empresa = int(df_abertas["onde_se_encontra"].isin(VALORES_ONDE_ENCONTRA_EMPRESA).sum())
    n_terceirizada = int(df_abertas["onde_se_encontra"].isin(VALORES_ONDE_ENCONTRA_TERCEIRIZADA).sum())
    n_fab = total - n_empresa - n_terceirizada

    caixa = slide.shapes.add_textbox(Emu(340000), topo + Emu(120000), Emu(largura_total - 680000), Emu(400000))
    run = caixa.text_frame.paragraphs[0].add_run()
    run.text = (
        f"{total} em aberto  ·  {n_empresa} com a empresa (VEE ONE)  ·  "
        f"{n_terceirizada} nas terceirizadas  ·  {n_fab} nas bases FAB"
    )
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COR_AMBER

    top_grafico = topo + Emu(460000)
    largura_grafico = Emu((largura_total - 3 * 300000) // 2)

    _legenda_grafico(slide, "Onde se encontra", largura_grafico, Emu(300000), top_grafico)
    onde_grupos = {
        "Com a empresa": n_empresa,
        "Terceirizadas": n_terceirizada,
        "Bases FAB": n_fab,
    }
    _grafico_barra(slide, list(onde_grupos.keys()), list(onde_grupos.values()), COR_GRAFICO_AMBER,
                   Emu(300000), top_grafico + Emu(320000), largura_grafico, Emu(4700000))

    left2 = Emu(300000) + largura_grafico + Emu(300000)
    _legenda_grafico(slide, "Distribuição por condição", largura_grafico, left2, top_grafico)
    condicao_contagem = df_abertas["condicao"].value_counts()
    _grafico_barra(slide, condicao_contagem.index, condicao_contagem.values, COR_GRAFICO_CYAN,
                   left2, top_grafico + Emu(320000), largura_grafico, Emu(4700000))


def gerar_apresentacao(ano, mes, caminho_saida, caminho_referencia=None):
    """Gera a apresentação do mês em `caminho_saida`, a partir da base
    `caminho_referencia` (padrão: 04_Relatorios/RMA_referencia.pptx)."""
    caminho_referencia = Path(caminho_referencia) if caminho_referencia else REFERENCIA
    if not caminho_referencia.exists():
        raise FileNotFoundError(
            f"Base da apresentação não encontrada em {caminho_referencia}. "
            "Baixe uma cópia da apresentação de um mês já pronto e salve nesse caminho."
        )

    aeronaves, resumo = _carregar_dados_mes(ano, mes)
    if not aeronaves:
        raise ValueError(f"Nenhuma aeronave negativou em {mes:02d}/{ano} — nada pra gerar nos slides por aeronave.")

    prs = Presentation(str(caminho_referencia))
    largura_total = prs.slide_width
    altura_total = prs.slide_height

    _atualizar_titulo_objetivo(prs, mes)

    # MMAM/P/PMAX/IFD oficiais — mesma planilha "RMA em andamento" que a
    # Ata de Reunião usa, não o resumo (local) do Cômputo Mensal (ver
    # docstring de `_construir_tabela_matriz`).
    arquivos_pasta_mes = _localizar_pasta_mes(ano, mes)
    conteudo_rma_oficial = _baixar_rma_em_andamento(arquivos_pasta_mes)
    indicadores_oficiais = extrair_indicadores_rma(conteudo_rma_oficial, ano, mes)

    caminho_saida_tmp = Path(caminho_saida)
    df_matriz_mes = pd.read_csv(DADOS_TRATADOS / "computo_mensal" / f"{ano}-{mes:02d}_matriz.csv", dtype={"matricula": str})
    for i, indice_matriz in enumerate(SLIDES_MATRIZ):
        _construir_tabela_matriz(
            prs.slides[indice_matriz], df_matriz_mes,
            resumo["aeronaves_pontuadas"], resumo["aeronaves_fora_listadas"],
            resumo, indicadores_oficiais, ano, mes, largura_total, altura_total - Emu(1300000),
            caminho_saida_tmp.parent / f"_matriz_{ano}-{mes:02d}_{i}.png",
        )

    df_pagamentos, contrato_resumo = _carregar_dados_pagamentos()
    _construir_slide_pagamentos(prs.slides[SLIDE_PAGAMENTOS], df_pagamentos, contrato_resumo, largura_total)

    df_emprestimos_mes, df_emprestimos_todos = _carregar_dados_emprestimos(ano, mes)
    reparaveis_abertas, reparaveis_filtradas = _carregar_dados_reparaveis()
    atrasos_abertas, atrasos_concluidas_mes, atrasos_resumo = _carregar_dados_atrasos(ano, mes)

    xml_slides = prs.slides._sldIdLst
    elemento_pagamentos = list(xml_slides)[SLIDE_PAGAMENTOS]
    elemento_emprestimos = list(xml_slides)[SLIDE_EMPRESTIMOS]
    elemento_dashboard_emprestimos = list(xml_slides)[SLIDE_EMPRESTIMOS_DASHBOARD]
    elemento_reparaveis_todas = list(xml_slides)[SLIDES_REPARAVEIS[0]]
    elemento_reparaveis_filtradas = list(xml_slides)[SLIDES_REPARAVEIS[1]]
    elemento_dashboard_atrasos = list(xml_slides)[SLIDE_ATRASOS_DASHBOARD]
    # Slides com conteúdo COPIADO da referência que nunca é regenerado por
    # nós — capturados AQUI (antes de qualquer inserção) pra remover no
    # final, mesmo padrão de identidade estável de elemento XML usado pros
    # outros slides antigos. Ver SLIDES_COPIADOS_EXTRAS e docstring do
    # módulo (pedido do Wallace, 2026-07-13).
    elementos_aeronave_antigos = [list(xml_slides)[i] for i in range(SLIDES_AERONAVE_INICIO, SLIDES_AERONAVE_FIM + 1)]
    elementos_copiados_extras = [list(xml_slides)[i] for i in SLIDES_COPIADOS_EXTRAS]

    # --- 1 slide por aeronave que negativou, inserido logo após o início
    # do bloco de aeronaves da referência ---
    indice_insercao = SLIDES_AERONAVE_INICIO - 1
    for matricula, dados_aeronave in aeronaves.items():
        slide = _slide_em_branco(prs, indice_insercao)
        _montar_slide_aeronave(slide, matricula, dados_aeronave, largura_total, ano, mes)
        indice_insercao += 1

    mes_nome = MESES_PT_MAIUSCULO[mes - 1].capitalize()

    indice_pagamentos_atual = list(xml_slides).index(elemento_pagamentos)
    slide_estatisticas_emprestimos = _slide_em_branco(prs, indice_pagamentos_atual)
    _construir_slide_emprestimos_estatisticas(slide_estatisticas_emprestimos, df_emprestimos_mes, f"Empréstimos — Estatísticas ({mes_nome})", largura_total)

    indice_emprestimos_atual = list(xml_slides).index(elemento_emprestimos)
    ultimo_indice_lista_emprestimos = _construir_paginas_tabela(
        prs, indice_emprestimos_atual - 1, df_emprestimos_mes, COLUNAS_EMPRESTIMO,
        f"Empréstimo em {mes_nome}", largura_total, Pt(5), 115000,
    )

    slide_estatisticas_emprestimos_total = _slide_em_branco(prs, ultimo_indice_lista_emprestimos)
    _construir_slide_emprestimos_estatisticas(slide_estatisticas_emprestimos_total, df_emprestimos_todos, "Empréstimos — Estatística Total (todos os meses)", largura_total)

    indice_reparaveis_todas_atual = list(xml_slides).index(elemento_reparaveis_todas)
    slide_estatisticas_reparaveis = _slide_em_branco(prs, indice_reparaveis_todas_atual - 1)
    _construir_slide_reparaveis_estatisticas(slide_estatisticas_reparaveis, reparaveis_abertas, largura_total)

    indice_dashboard_atrasos_atual = list(xml_slides).index(elemento_dashboard_atrasos)
    slide_atrasos_situacao = _slide_em_branco(prs, indice_dashboard_atrasos_atual - 1)
    _construir_slide_atrasos_situacao(slide_atrasos_situacao, atrasos_abertas, largura_total)

    # slide_atrasos_situacao acabou de ser inserido em
    # (indice_dashboard_atrasos_atual - 1) + 1 = indice_dashboard_atrasos_atual —
    # inserir o próximo logo depois dele, sem precisar relocalizar via objeto.
    slide_atrasos_entregas = _slide_em_branco(prs, indice_dashboard_atrasos_atual)
    _construir_slide_atrasos_entregas(slide_atrasos_entregas, atrasos_concluidas_mes, atrasos_resumo, largura_total)

    elementos_remover = (
        [elemento_emprestimos, elemento_dashboard_emprestimos, elemento_reparaveis_todas,
         elemento_reparaveis_filtradas, elemento_dashboard_atrasos]
        + elementos_aeronave_antigos + elementos_copiados_extras
    )
    for elemento in elementos_remover:
        indice_atual = list(xml_slides).index(elemento)
        _remover_slides(prs, indice_atual, indice_atual)

    caminho_saida = Path(caminho_saida)
    caminho_saida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(caminho_saida))
    return caminho_saida


if __name__ == "__main__":
    import sys
    ano_arg = int(sys.argv[1])
    mes_arg = int(sys.argv[2])
    destino = DADOS_TRATADOS / "atas" / f"RMA_{MESES_PT_MAIUSCULO[mes_arg-1].capitalize()}_{ano_arg}_TESTE.pptx"
    caminho = gerar_apresentacao(ano_arg, mes_arg, destino)
    print(f"Gerado: {caminho}")
